import os
from PyPDF2 import PdfReader
from docx import Document
from pptx import Presentation
from rag_system import RAGSystem

def extract_text_from_pdf(file_path):
    reader = PdfReader(file_path)
    return "\n".join([page.extract_text() for page in reader.pages if page.extract_text()])

def extract_text_from_docx(file_path):
    doc = Document(file_path)
    return "\n".join([para.text for para in doc.paragraphs])

def extract_text_from_pptx(file_path):
    prs = Presentation(file_path)
    text_runs = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text_runs.append(shape.text)
    return "\n".join(text_runs)

def extract_text_from_txt(file_path):
    with open(file_path, 'r', encoding='utf-8') as f:
        return f.read()

def universal_extractor(file_path):
    """
    Routes the file to the correct extractor based on extension.
    """
    ext = os.path.splitext(file_path)[1].lower()
    
    try:
        if ext == '.pdf':
            return extract_text_from_pdf(file_path)
        elif ext == '.docx':
            return extract_text_from_docx(file_path)
        elif ext == '.pptx':
            return extract_text_from_pptx(file_path)
        elif ext == '.txt':
            return extract_text_from_txt(file_path)
        else:
            return f"Unsupported file format: {ext}"
    except Exception as e:
        return f"Error processing {file_path}: {str(e)}"

def chunk_text(text, chunk_size=1000, overlap=200):
    """
    Splits text into chunks of `chunk_size` characters with `overlap` characters of overlap.
    """
    if not text:
        return []
    
    chunks = []
    start = 0
    text_len = len(text)
    
    while start < text_len:
        end = start + chunk_size
        chunks.append(text[start:end])
        if end >= text_len:
            break
        start += chunk_size - overlap
        
    return chunks

def process_and_store_document(file_path, rag_system=None,collection_name=None):
    """
    Extracts text from a document, chunks it, and stores the chunks in Elasticsearch.
    """
    print(f"Extracting text from: {file_path}")
    text = universal_extractor(file_path)
    
    if text.startswith("Unsupported file format:") or text.startswith("Error processing "):
        print(text)
        return False
        
    print("Chunking extracted text...")
    chunks = chunk_text(text, chunk_size=1000, overlap=200)
    print(f"Generated {len(chunks)} chunks.")
    
    if not chunks:
        print("No text found to process.")
        return False
        
    if rag_system is None:
        print("Initializing RAG System...")
        rag_system = RAGSystem(index_name=collection_name)
        
    metadatas = [
        {"source_file": os.path.basename(file_path), "chunk_id": i} 
        for i in range(len(chunks))
    ]
    
    rag_system.add_documents(chunks, metadatas=metadatas)
    print(f"Successfully processed and stored {file_path}")
    return True

# if __name__ == "__main__":
#     # Example Usage:
#     process_and_store_document("bajaj_life.txt")
#     pass